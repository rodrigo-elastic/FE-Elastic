"""
filename: routes_weekly_slides.py
description: Weekly customer status slide deck. Aggregates post-meeting records
for a given week, groups by company, and uses Claude to synthesize slide content
matching the Field Engineering weekly standup format (Actions, Renewals, Cases,
Consumption, Feature Adoption, Risks/Notes/Top of mind).
date: 09-05-2026
"""
__author__ = "Rodrigo Careaga"
__copyright__ = "Copyright 2026, Rodrigo Careaga"
__version__ = "0.1.0"
__status__ = "Development"

import json
from datetime import date, datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional

import httpx
from fastapi import APIRouter, HTTPException, Query
from pydantic import BaseModel, Field

from app.config import settings
from app.integrations.claude_client import get_elastic_service
from app.utils.logging import get_logger

log = get_logger(__name__)

router = APIRouter(prefix="/weekly-slides", tags=["weekly-slides"])


# ============================================================ Pydantic output schema =


class SlideRenewal(BaseModel):
    label: str = ""
    notes: str = ""
    risk: str = ""
    amount: str = ""
    date: str = ""


class WeeklySlideOut(BaseModel):
    use_case: str = ""
    temperature: str = "stable"
    temperature_reason: str = ""
    current_actions: List[str] = Field(default_factory=list)
    upcoming_actions: List[str] = Field(default_factory=list)
    renewals: List[SlideRenewal] = Field(default_factory=list)
    cases: List[str] = Field(default_factory=list)
    consumption: str = ""
    wow_pct: str = "N/A"
    feature_adoption: List[str] = Field(default_factory=list)
    risks_notes: List[str] = Field(default_factory=list)


_SLIDE_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "use_case": {"type": "string"},
        "temperature": {"type": "string", "enum": ["churn", "stable", "growth"]},
        "temperature_reason": {"type": "string"},
        "current_actions": {"type": "array", "items": {"type": "string"}},
        "upcoming_actions": {"type": "array", "items": {"type": "string"}},
        "renewals": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "label": {"type": "string"},
                    "notes": {"type": "string"},
                    "risk": {"type": "string"},
                    "amount": {"type": "string"},
                    "date": {"type": "string"},
                },
                "required": ["label"],
            },
        },
        "cases": {"type": "array", "items": {"type": "string"}},
        "consumption": {"type": "string"},
        "wow_pct": {"type": "string"},
        "feature_adoption": {"type": "array", "items": {"type": "string"}},
        "risks_notes": {"type": "array", "items": {"type": "string"}},
    },
    "required": [
        "use_case", "temperature", "current_actions",
        "upcoming_actions", "consumption", "risks_notes",
    ],
}


# ============================================================ Helpers ===============


def _week_bounds(week_start_str: Optional[str]):
    if week_start_str:
        try:
            d = date.fromisoformat(week_start_str)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid week_start. Use YYYY-MM-DD.")
    else:
        today = date.today()
        d = today - timedelta(days=today.weekday())
    start = datetime(d.year, d.month, d.day, tzinfo=timezone.utc)
    return start, start + timedelta(days=7)


def _load_post_meetings(demo_mode: bool, start: datetime, end: datetime) -> List[Dict[str, Any]]:
    records: List[Dict[str, Any]] = []
    seen: set = set()

    try:
        from app.repositories.elasticsearch_repo import get_repo as get_es_repo
        es = get_es_repo()
        if es.available:
            for rec in es.list_post_meetings(limit=500):
                mid = rec.get("meeting_id", "")
                if mid in seen:
                    continue
                if not demo_mode:
                    ga = rec.get("generated_at", "")
                    if not ga:
                        continue
                    try:
                        ts = datetime.fromisoformat(ga.replace("Z", "+00:00"))
                        if not (start <= ts < end):
                            continue
                    except Exception:
                        continue
                records.append(rec)
                seen.add(mid)
    except Exception as exc:
        log.warning("weekly_slides.es_load_failed", error=str(exc))

    post_dir = settings.runtime_dir / "post_meeting"
    if post_dir.exists():
        for p in sorted(post_dir.glob("*.json")):
            try:
                rec = json.loads(p.read_text(encoding="utf-8"))
                mid = rec.get("meeting_id", p.stem)
                if mid in seen:
                    continue
                if not demo_mode:
                    ga = rec.get("generated_at", "")
                    if not ga:
                        continue
                    ts = datetime.fromisoformat(ga.replace("Z", "+00:00"))
                    if not (start <= ts < end):
                        continue
                records.append(rec)
                seen.add(mid)
            except Exception:
                pass

    return records


def _group_by_company(records: List[Dict[str, Any]]) -> Dict[str, List[Dict[str, Any]]]:
    grouped: Dict[str, List[Dict[str, Any]]] = {}
    for rec in records:
        key = (rec.get("company_name") or rec.get("company_id") or "Unknown").strip()
        grouped.setdefault(key, []).append(rec)
    return grouped


def _extract_sf(records: List[Dict[str, Any]]) -> Dict[str, Any]:
    for rec in sorted(records, key=lambda r: r.get("generated_at", ""), reverse=True):
        sw = rec.get("salesforce_writes") or {}
        if sw:
            return sw
    return {}


def _mock_slide(company_name: str) -> Dict[str, Any]:
    """Return a demo slide payload without calling any LLM."""
    return {
        "company_name": company_name,
        "use_case": "Security + Observability",
        "temperature": "stable",
        "temperature_reason": "Active engagement. No major risk signals.",
        "current_actions": [
            f"Follow up on open items from {company_name}",
            "Send meeting summary and next steps",
        ],
        "upcoming_actions": ["Schedule next business review"],
        "renewals": [],
        "cases": [],
        "consumption": "Stable. On-prem limits telemetry visibility.",
        "wow_pct": "N/A",
        "feature_adoption": ["Elastic Stack", "Observability"],
        "risks_notes": ["Review pending action item ownership"],
        "arr": "",
        "cloud_arr": "",
        "training_services": "",
        "renewable_base": "",
        "open_ne": "",
        "salesforce_url": "",
        "meeting_count": 1,
        "meeting_ids": [],
        "updated": date.today().isoformat(),
    }


def _build_slide(company_name: str, records: List[Dict[str, Any]], sf: Dict[str, Any]) -> Dict[str, Any]:
    summaries = [r.get("summary", "") for r in records if r.get("summary")]
    all_actions: List[Dict[str, Any]] = []
    for rec in records:
        all_actions.extend(rec.get("action_items") or [])
    meddpicc: List[Dict[str, Any]] = []
    for rec in records:
        meddpicc.extend(rec.get("meddpicc_signals") or [])

    opp = sf.get("opportunity") or {}
    account = sf.get("account") or {}

    sf_ctx_lines = []
    if account.get("Name"):
        sf_ctx_lines.append(f"Account: {account['Name']}, Industry: {account.get('Industry', 'N/A')}")
    if opp.get("Name"):
        amt = opp.get("Amount") or 0
        sf_ctx_lines.append(
            f"Opportunity: {opp['Name']}, Stage: {opp.get('StageName', '')}, "
            f"Amount: ${amt:,}, Close: {opp.get('CloseDate', '')}"
        )
    sf_ctx = "\n".join(sf_ctx_lines) or "Not available."

    ai_lines = "\n".join(
        f"- {ai.get('title', '')} (owner: {ai.get('owner_name', 'TBD')}, "
        f"due: {ai.get('due_date', 'TBD')}, impact: {ai.get('impact', 'med')})"
        for ai in all_actions
    ) or "No action items."

    meddpicc_lines = "\n".join(
        f"- [{m.get('category', '')}] {m.get('note', '') or m.get('quote', '')[:120]}"
        for m in meddpicc[:8]
    ) or "No MEDDPICC signals."

    summary_text = "\n\n---\n\n".join(summaries) or "No meeting summaries available."

    system = (
        "You are a Field Engineering weekly standup assistant at Elastic. "
        "You produce structured, concise, executive-ready customer status summaries. "
        "Be specific - use real names, amounts, dates from the input data. "
        "Never invent facts not in the input. If data is missing, say so briefly."
    )

    user = f"""Generate a weekly customer status slide for the FE team standup.

Company: {company_name}
Salesforce: {sf_ctx}

Meeting summaries this week:
{summary_text}

Action items:
{ai_lines}

MEDDPICC signals:
{meddpicc_lines}

Return a JSON object with these exact keys:
{{
  "use_case": "Short description of main Elastic use cases (e.g. 'Security + Observability')",
  "temperature": "churn|stable|growth",
  "temperature_reason": "One sentence explaining the account health signal.",
  "current_actions": ["2-5 in-flight action items. Include owner. One concise line each."],
  "upcoming_actions": ["1-3 planned future actions."],
  "renewals": [{{"label": "Renewal name/description", "notes": "Brief context", "risk": "low|med|high", "amount": "$amount", "date": "YYYY-MM-DD"}}],
  "cases": ["L2 - Title - Status (use format from data if available)"],
  "consumption": "1-2 sentence summary of consumption trend and account health.",
  "wow_pct": "??% or actual week-over-week % if inferable from data",
  "feature_adoption": ["Elastic feature being actively adopted - max 4"],
  "risks_notes": ["Key risk, note, or top-of-mind item - 2-4 items"]
}}

Rules:
- temperature: 'growth' if positive deal signals or expansion; 'churn' if risk, disengagement, or renewal concerns; 'stable' otherwise.
- current_actions: in-flight or due soon. upcoming_actions: planned for future meetings/sprints.
- cases: [] if no case data available.
- renewals: use opportunity data; [] if none.
- Respond with ONLY the JSON object. No markdown. No explanation."""

    mock_payload: Dict[str, Any] = {
        "use_case": "Security + Observability",
        "temperature": "stable",
        "temperature_reason": "Active engagement across multiple workstreams. No major risk signals.",
        "current_actions": [
            f"Follow up on open action items from {company_name} meetings",
            "Send meeting summary and next steps",
        ],
        "upcoming_actions": ["Schedule next business review"],
        "renewals": [],
        "cases": [],
        "consumption": "Stable consumption. On-prem deployment limits telemetry visibility.",
        "wow_pct": "N/A",
        "feature_adoption": ["Elastic Stack", "Agent Builder", "Observability"],
        "risks_notes": ["Review pending action items and ownership assignments"],
    }

    model_name = settings.model_for("post_meeting")

    # Always use the Elastic inference connector - customer data must not leave
    # the Elastic infrastructure. strict=True blocks all fallback paths to the
    # direct Anthropic API.
    try:
        svc = get_elastic_service()
        result: WeeklySlideOut = svc.call_structured(
            system=system,
            user=user,
            schema=_SLIDE_SCHEMA,
            output_model=WeeklySlideOut,
            model=model_name,
            max_tokens=1500,
            effort="high",
            thinking_adaptive=True,
            cache_system=True,
            mock_payload=mock_payload,
            audit_meta={"agent": "weekly_slides", "company": company_name},
            strict=True,
        )
        slide = result.model_dump()
    except RuntimeError as exc:
        # get_elastic_service() raised - Kibana not configured or strict block triggered.
        log.warning("weekly_slides.elastic_required", company=company_name, error=str(exc)[:200])
        raise
    except Exception as exc:
        log.warning("weekly_slides.claude_failed", company=company_name, error=str(exc)[:200])
        slide = mock_payload.copy()

    arr = opp.get("Amount") or 0
    slide["company_name"] = company_name
    slide["arr"] = f"${arr:,}" if arr else ""
    slide["cloud_arr"] = ""
    slide["training_services"] = ""
    slide["renewable_base"] = ""
    slide["open_ne"] = ""
    slide["salesforce_url"] = account.get("Url") or opp.get("Url") or ""
    slide["meeting_count"] = len(records)
    slide["meeting_ids"] = [r.get("meeting_id", "") for r in records]
    slide["updated"] = date.today().isoformat()
    return slide


# ============================================================ PPTX builder ==========


def _build_pptx(slide_data: Dict[str, Any]) -> "Path":
    """Generate a .pptx matching the FE standup slide template exactly.

    Layout (top to bottom):
      HEADER : logo | ARR pills | company name / use case / date | meta pills | temperature box
      BODY   : Current actions (wide) | Renewals | Cases | Consumption
      BOTTOM : Feature Adoption | Risks / Notes / Top of mind
      FOOTER : icon labels row
    """
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.dml.color import RGBColor  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]
        from pptx.util import Inches, Pt  # type: ignore[import]
    except ImportError as exc:
        raise RuntimeError("python-pptx required: pip install python-pptx") from exc

    # ── Palette ────────────────────────────────────────────────────
    RED   = RGBColor(0xE8, 0x4B, 0x37)  # actions / feature adoption header
    NAVY  = RGBColor(0x0F, 0x2D, 0x5C)  # renewals / temperature header
    ORNGE = RGBColor(0xEB, 0x6B, 0x4A)  # cases header
    TEAL  = RGBColor(0x00, 0xB4, 0xA2)  # consumption header
    GOLD  = RGBColor(0xF1, 0xA7, 0x30)  # risks/notes header  (also Stable pill)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY = RGBColor(0xF0, 0xF0, 0xF0)  # section body background
    DGRAY = RGBColor(0xCC, 0xCC, 0xCC)  # pills / logo border
    DTEXT = RGBColor(0x1A, 0x1A, 0x1A)
    MUTED = RGBColor(0x88, 0x88, 0x88)

    C_CHURN_ON   = RGBColor(0xE8, 0x4B, 0x37)
    C_CHURN_OFF  = RGBColor(0xF5, 0xC4, 0xBE)
    C_STABLE_ON  = RGBColor(0xF1, 0xA7, 0x30)
    C_STABLE_OFF = RGBColor(0xF8, 0xE8, 0xC5)
    C_GROWTH_ON  = RGBColor(0x3C, 0xB4, 0x4B)
    C_GROWTH_OFF = RGBColor(0xC5, 0xEC, 0xCC)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    I = Inches

    # ── Low-level helpers ──────────────────────────────────────────

    def _rect(l, t, w, h, bg, border=None):
        s = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
        s.fill.solid()
        s.fill.fore_color.rgb = bg
        if border:
            s.line.color.rgb = border
            s.line.width = Pt(0.75)
        else:
            s.line.fill.background()
        return s

    def _pill(l, t, w, h, text, bg, fg=WHITE, bold=True, size=8.5, align=PP_ALIGN.CENTER):
        s = slide.shapes.add_shape(5, I(l), I(t), I(w), I(h))
        s.fill.solid()
        s.fill.fore_color.rgb = bg
        s.line.fill.background()
        tf = s.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(text)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = fg

    def _new_tf(l, t, w, h, wrap=True):
        box = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        return tf

    def _para(tf, text, size=7.5, bold=False, italic=False,
               color=None, align=PP_ALIGN.LEFT, space_before=0):
        """Append one paragraph to tf, reusing the first empty para."""
        if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(space_before)
        r = p.add_run()
        r.text = str(text)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color if color else DTEXT
        return p

    def _hdr(l, t, w, title, bg, underline=False, size=9.5):
        """Colored section header bar. Returns bar height (inches)."""
        H = 0.30
        s = slide.shapes.add_shape(1, I(l), I(t), I(w), I(H))
        s.fill.solid()
        s.fill.fore_color.rgb = bg
        s.line.fill.background()
        p = s.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.underline = underline
        r.font.color.rgb = WHITE
        return H

    # ── Grid constants ─────────────────────────────────────────────
    # Matching the template proportions from the images:
    #   actions ~33% | renewals ~16% | cases ~26% | consumption ~25%
    C1_W = 4.44   # Current and upcoming actions
    C2_W = 2.10   # Renewals
    C3_W = 3.48   # Cases
    C4_W = 13.33 - C1_W - C2_W - C3_W  # Consumption  = 3.31

    C1_X, C2_X, C3_X = 0.0, C1_W, C1_W + C2_W
    C4_X = C3_X + C3_W

    H_HDR  = 1.48   # header zone height
    BODY_T = H_HDR
    BODY_H = 3.88
    BTM_T  = BODY_T + BODY_H   # ~5.36
    BTM_H  = 1.30
    FOOT_T = BTM_T + BTM_H     # ~6.66
    FOOT_H = 7.5 - FOOT_T      # ~0.84

    HDR_H  = 0.30   # section header bar height

    # Bottom row
    BTM_L_W = C1_W            # Feature Adoption (same as actions column)
    BTM_R_W = 13.33 - BTM_L_W # Risks / Notes / Top of mind

    # Temperature box (far right of header)
    TEMP_W = 2.72
    TEMP_X = 13.33 - TEMP_W

    # Meta pills zone (just left of temperature)
    META_W = 1.92
    META_X = TEMP_X - META_W - 0.07

    # Company name zone (center)
    CO_X = 3.78
    CO_W = META_X - CO_X - 0.08

    # ── HEADER ─────────────────────────────────────────────────────

    # Logo box (dashed border)
    _rect(0.05, 0.06, 1.62, H_HDR - 0.08, WHITE, border=DGRAY)
    tf = _new_tf(0.05, 0.52, 1.62, 0.40)
    _para(tf, "Company logo", size=8, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # ARR pills (left of company name)
    ARR_X, ARR_W, ARR_H = 1.78, 1.86, 0.30
    arr_val  = slide_data.get("arr")       or ""
    carr_val = slide_data.get("cloud_arr") or ""

    _rect(ARR_X, 0.14, ARR_W, ARR_H, DGRAY)
    tf = _new_tf(ARR_X + 0.08, 0.16, ARR_W - 0.16, ARR_H - 0.04)
    _para(tf, "Total ARR:", size=6.5, bold=True)
    _para(tf, arr_val or "-", size=7.5, bold=True)

    _rect(ARR_X, 0.14 + ARR_H + 0.07, ARR_W, ARR_H, DGRAY)
    tf = _new_tf(ARR_X + 0.08, 0.14 + ARR_H + 0.09, ARR_W - 0.16, ARR_H - 0.04)
    _para(tf, "Cloud ARR:", size=6.5, bold=True)
    _para(tf, carr_val or "-", size=7.5, bold=True)

    # Company name (large bold) + use case + updated date
    co_name = slide_data.get("company_name") or ""
    co_use  = slide_data.get("use_case")     or ""
    co_date = slide_data.get("updated")      or ""

    tf = _new_tf(CO_X, 0.04, CO_W, 0.78)
    _para(tf, co_name, size=24, bold=True, align=PP_ALIGN.CENTER)

    tf = _new_tf(CO_X, 0.82, CO_W, 0.36)
    _para(tf, co_use, size=13, align=PP_ALIGN.CENTER)

    tf = _new_tf(CO_X, 1.17, CO_W, 0.28)
    _para(tf, f"Updated: {co_date}", size=8, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # Meta pills: Training/Services, Renewable Base, Open N&E
    META_H, META_GAP = 0.28, 0.08
    meta_rows = [
        ("Training/Services:", slide_data.get("training_services") or ""),
        ("Renewable Base:",    slide_data.get("renewable_base")    or ""),
        ("Open N&E:",          slide_data.get("open_ne")           or ""),
    ]
    for i, (lbl, val) in enumerate(meta_rows):
        my = 0.13 + i * (META_H + META_GAP)
        _rect(META_X, my, META_W, META_H, DGRAY)
        tf = _new_tf(META_X + 0.07, my + 0.04, META_W - 0.14, META_H - 0.06)
        _para(tf, f"{lbl}  {val}" if val else lbl, size=7, color=DTEXT)

    # Temperature box background
    _rect(TEMP_X, 0, TEMP_W, H_HDR, LGRAY)

    # "Potential / Temperature" blue header bar
    s = slide.shapes.add_shape(1, I(TEMP_X), I(0), I(TEMP_W), I(0.36))
    s.fill.solid()
    s.fill.fore_color.rgb = NAVY
    s.line.fill.background()
    p = s.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Potential / Temperature"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # Arrow
    tf = _new_tf(TEMP_X, 0.37, TEMP_W, 0.30, wrap=False)
    _para(tf, "▼", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Churn / Stable / Growth pills
    temp = (slide_data.get("temperature") or "stable").lower()
    PILL_W, PILL_H, PILL_Y = 0.80, 0.33, 0.72
    pill_gap = (TEMP_W - PILL_W * 3) / 4
    for pi, (lbl, on_c, off_c) in enumerate([
        ("Churn",  C_CHURN_ON,  C_CHURN_OFF),
        ("Stable", C_STABLE_ON, C_STABLE_OFF),
        ("Growth", C_GROWTH_ON, C_GROWTH_OFF),
    ]):
        active = temp == lbl.lower()
        _pill(TEMP_X + pill_gap + pi * (PILL_W + pill_gap), PILL_Y,
              PILL_W, PILL_H, lbl,
              on_c if active else off_c,
              fg=WHITE if active else DTEXT,
              bold=True, size=9.5)

    reason = (slide_data.get("temperature_reason") or "")[:100]
    if reason:
        tf = _new_tf(TEMP_X + 0.06, 1.10, TEMP_W - 0.12, 0.36)
        _para(tf, reason, size=6.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # ── BODY ROW ─────────────────────────────────────────────────
    # Col 1 - Current and upcoming actions (red)
    _hdr(C1_X, BODY_T, C1_W, "Current and upcoming actions", RED)
    _rect(C1_X, BODY_T + HDR_H, C1_W, BODY_H - HDR_H, LGRAY)

    current  = slide_data.get("current_actions")  or []
    upcoming = slide_data.get("upcoming_actions") or []
    tf = _new_tf(C1_X + 0.12, BODY_T + HDR_H + 0.12, C1_W - 0.24, BODY_H - HDR_H - 0.18)
    if current:
        _para(tf, "Current Actions", size=8.5, bold=True, align=PP_ALIGN.CENTER)
        for a in current[:5]:
            _para(tf, a, size=7.5, space_before=3)
    if upcoming:
        _para(tf, "Upcoming actions", size=8.5, bold=True,
              align=PP_ALIGN.CENTER, space_before=10 if current else 0)
        for a in upcoming[:4]:
            _para(tf, a, size=7.5, space_before=3)
    if not current and not upcoming:
        _para(tf, "No actions recorded.", size=7.5, color=MUTED)

    # Col 2 - Renewals (navy)
    _hdr(C2_X, BODY_T, C2_W, "Renewals", NAVY)
    _rect(C2_X, BODY_T + HDR_H, C2_W, BODY_H - HDR_H, LGRAY)

    renewals = slide_data.get("renewals") or []
    tf = _new_tf(C2_X + 0.09, BODY_T + HDR_H + 0.10, C2_W - 0.18, BODY_H - HDR_H - 0.15)
    if renewals:
        for r in renewals[:4]:
            lbl   = r.get("label") or ""
            amt   = r.get("amount") or ""
            dt    = r.get("date")   or ""
            notes = r.get("notes")  or ""
            risk  = r.get("risk")   or ""
            line  = f"- {lbl}"
            extras = " - ".join(x for x in [amt, dt] if x)
            if extras:
                line += f" - {extras}"
            _para(tf, line, size=7.5, space_before=5)
            if notes:
                _para(tf, f"  - Notes: {notes}", size=7, space_before=1)
            if risk:
                _para(tf, f"  - Risk: {risk}", size=7, space_before=1)
    else:
        _para(tf, "No renewals on record.", size=7.5, color=MUTED)

    # Col 3 - Cases (orange/salmon)
    _hdr(C3_X, BODY_T, C3_W, "Cases", ORNGE)
    _rect(C3_X, BODY_T + HDR_H, C3_W, BODY_H - HDR_H, LGRAY)

    cases = slide_data.get("cases") or []
    tf = _new_tf(C3_X + 0.09, BODY_T + HDR_H + 0.10, C3_W - 0.18, BODY_H - HDR_H - 0.15)
    if cases:
        for c in cases[:6]:
            _para(tf, f"- {c}", size=7.5, space_before=5)
    else:
        _para(tf, "No open cases.", size=7.5, color=MUTED)

    # Col 4 - Consumption (teal, underlined title)
    _hdr(C4_X, BODY_T, C4_W, "Consumption", TEAL, underline=True)
    _rect(C4_X, BODY_T + HDR_H, C4_W, BODY_H - HDR_H, LGRAY)

    cons = (slide_data.get("consumption") or "")[:320]
    wow  = slide_data.get("wow_pct") or "N/A"

    # "WoW:" label then gray pill for the percentage
    tf = _new_tf(C4_X + 0.09, BODY_T + HDR_H + 0.10, 0.55, 0.28)
    _para(tf, "WoW:", size=8.5, bold=True)
    _pill(C4_X + 0.66, BODY_T + HDR_H + 0.09, 0.85, 0.28, wow, DGRAY,
          fg=DTEXT, bold=False, size=8)

    tf = _new_tf(C4_X + 0.09, BODY_T + HDR_H + 0.46, C4_W - 0.18, BODY_H - HDR_H - 0.55)
    _para(tf, cons, size=7.5, italic=True)

    # ── BOTTOM ROW ─────────────────────────────────────────────────
    # Left: Feature Adoption (red)
    _hdr(0, BTM_T, BTM_L_W, "Feature Adoption", RED)
    _rect(0, BTM_T + HDR_H, BTM_L_W, BTM_H - HDR_H, LGRAY)

    features = slide_data.get("feature_adoption") or []
    tf = _new_tf(0.12, BTM_T + HDR_H + 0.10, BTM_L_W - 0.24, BTM_H - HDR_H - 0.14)
    if features:
        for i, f in enumerate(features[:4]):
            _para(tf, f"Feature {i + 1}: {f}", size=7.5, space_before=3)
    else:
        _para(tf, "No feature data.", size=7.5, color=MUTED)

    # Right: Risks / Notes / Top of mind (gold)
    _hdr(BTM_L_W, BTM_T, BTM_R_W, "Risks/ Notes / Top of mind", GOLD)
    _rect(BTM_L_W, BTM_T + HDR_H, BTM_R_W, BTM_H - HDR_H, LGRAY)

    risks = slide_data.get("risks_notes") or []
    tf = _new_tf(BTM_L_W + 0.12, BTM_T + HDR_H + 0.10, BTM_R_W - 0.24, BTM_H - HDR_H - 0.14)
    if risks:
        for r in risks[:4]:
            _para(tf, f"- {r}", size=7.5, space_before=3)
    else:
        _para(tf, "None noted.", size=7.5, color=MUTED)

    # ── FOOTER ─────────────────────────────────────────────────────
    _rect(0, FOOT_T, 13.33, FOOT_H, LGRAY)

    icon_labels = [
        "[Salesforce]", "[Consumption]", "[Contacts]",
        "[LinkedIn]", "[org. chart]", "[GDrive]", "elastic",
    ]
    cell_w = 13.33 / len(icon_labels)
    for i, lbl in enumerate(icon_labels):
        tf = _new_tf(cell_w * i, FOOT_T + 0.08, cell_w, FOOT_H - 0.16, wrap=False)
        _para(tf, lbl, size=7.5, color=MUTED, align=PP_ALIGN.CENTER)

    # ── Save ───────────────────────────────────────────────────────
    slides_dir = settings.runtime_dir / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)
    slug = "".join(c if c.isalnum() or c in "-_" else "_"
                   for c in (slide_data.get("company_name") or "slide"))
    updated = (slide_data.get("updated") or date.today().isoformat()).replace("-", "")[:8]
    out = slides_dir / f"{slug}_{updated}.pptx"
    prs.save(str(out))
    log.info("weekly_slides.pptx_saved", path=str(out))
    return out


# ============================================================ Slack upload ==========


def _slack_send_slide(pptx_path: "Path", company_name: str, channel: str, pptx_url: str) -> Dict[str, Any]:
    """Upload PPTX to Slack via bot token, or fall back to webhook message with link."""
    bot_token = settings.slack_bot_token.strip()
    webhook   = settings.slack_webhook_url.strip()

    # ── Bot token path: actual file upload ──────────────────────────
    if bot_token:
        try:
            with pptx_path.open("rb") as fh:
                resp = httpx.post(
                    "https://slack.com/api/files.upload",
                    headers={"Authorization": f"Bearer {bot_token}"},
                    data={
                        "channels": channel,
                        "filename": pptx_path.name,
                        "initial_comment": f":bar_chart: *{company_name}* - Weekly status slide",
                    },
                    files={"file": (pptx_path.name, fh,
                                    "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                    timeout=30.0,
                )
            data = resp.json()
            if data.get("ok"):
                log.info("weekly_slides.slack_uploaded", channel=channel, file=pptx_path.name)
                return {"ok": True, "method": "upload", "channel": channel,
                        "slack_url": data.get("file", {}).get("permalink", "")}
            log.warning("weekly_slides.slack_upload_failed", error=data.get("error"))
        except Exception as exc:
            log.warning("weekly_slides.slack_upload_error", error=str(exc)[:200])

    # ── Webhook fallback: post message with download link ───────────
    if webhook:
        try:
            msg = {
                "text": f":bar_chart: *{company_name}* - Weekly status slide ready",
                "attachments": [{
                    "color": "#00B4A2",
                    "text": f"<{pptx_url}|Download PPTX> | _{pptx_path.name}_",
                }],
            }
            resp = httpx.post(webhook, json=msg, timeout=10.0)
            if resp.status_code == 200:
                log.info("weekly_slides.slack_webhook_sent", channel=channel)
                return {"ok": True, "method": "webhook", "channel": channel, "slack_url": ""}
        except Exception as exc:
            log.warning("weekly_slides.slack_webhook_error", error=str(exc)[:200])

    return {"ok": False, "method": "none", "channel": channel, "error": "No Slack credentials configured."}


# ============================================================ Endpoints =============


@router.get("/download/{filename}")
def download_slide(filename: str):
    """Serve a generated PPTX file."""
    from fastapi.responses import FileResponse

    safe = "".join(c for c in filename if c.isalnum() or c in "-_.")
    path = settings.runtime_dir / "slides" / safe
    if not path.exists() or path.suffix != ".pptx":
        raise HTTPException(status_code=404, detail="Slide not found.")
    return FileResponse(
        str(path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=safe,
    )


@router.post("/from-meeting/{meeting_id}")
def slide_from_meeting(
    meeting_id: str,
    demo: bool = Query(False, description="Use mock slide content; skip LLM call"),
) -> Dict[str, Any]:
    """Generate a PPTX customer status slide from a single post-meeting record
    and send it to the configured Slack channel.

    Uses the Elastic inference connector (strict mode - customer data never
    leaves Elastic infrastructure). Pass demo=true to use mock content.
    """
    # ── Load the post-meeting record ────────────────────────────────
    rec: Optional[Dict[str, Any]] = None

    try:
        from app.repositories.elasticsearch_repo import get_repo as get_es_repo
        es = get_es_repo()
        if es.available:
            rec = es.get_post_meeting(meeting_id)
    except Exception:
        pass

    if rec is None:
        post_path = settings.runtime_dir / "post_meeting" / f"{meeting_id}.json"
        if post_path.exists():
            try:
                rec = json.loads(post_path.read_text(encoding="utf-8"))
            except Exception as exc:
                raise HTTPException(status_code=500, detail=f"Failed to read post-meeting: {exc}")

    if rec is None:
        raise HTTPException(status_code=404, detail=f"No post-meeting record found for meeting {meeting_id}.")

    company_name = (rec.get("company_name") or rec.get("company_id") or "Unknown").strip()
    sf = rec.get("salesforce_writes") or {}

    # ── Generate slide content via Elastic Claude ───────────────────
    if demo:
        slide_data = _mock_slide(company_name)
    else:
        try:
            slide_data = _build_slide(company_name, [rec], sf)
        except RuntimeError as exc:
            raise HTTPException(
                status_code=503,
                detail=f"Elastic inference connector required. {exc}",
            )

    # ── Build PPTX ─────────────────────────────────────────────────
    try:
        pptx_path = _build_pptx(slide_data)
    except Exception as exc:
        log.warning("weekly_slides.pptx_failed", meeting_id=meeting_id, error=str(exc)[:200])
        raise HTTPException(status_code=500, detail=f"PPTX generation failed: {exc}")

    # ── Resolve download URL ────────────────────────────────────────
    from app.api.routes_workflow_settings import _load_settings as _load_wf_settings
    try:
        wf = _load_wf_settings()
        slack_channel = wf.get("slack_channel") or "#fe-copilot-briefs"
    except Exception:
        slack_channel = "#fe-copilot-briefs"

    pptx_rel  = f"/api/v1/weekly-slides/download/{pptx_path.name}"
    pptx_url  = f"{settings.public_base_url}{pptx_rel}"

    # ── Send to Slack ───────────────────────────────────────────────
    slack_result = _slack_send_slide(pptx_path, company_name, slack_channel, pptx_url)

    log.info("weekly_slides.from_meeting_done",
             meeting_id=meeting_id, company=company_name,
             slack_ok=slack_result.get("ok"))

    return {
        "ok": True,
        "company_name": company_name,
        "meeting_id": meeting_id,
        "slide_name": pptx_path.name,
        "pptx_url": pptx_url,
        "pptx_rel": pptx_rel,
        "slack": slack_result,
        "slack_channel": slack_channel,
    }


@router.get("")
def get_weekly_slides(
    week_start: Optional[str] = Query(None, description="Monday YYYY-MM-DD. Defaults to current week."),
    demo: bool = Query(False, description="Use all available post-meetings regardless of date."),
) -> Dict[str, Any]:
    """Generate weekly customer status slides from post-meeting records.

    Groups meetings by company, calls Claude to synthesize slide content
    (actions, renewals, cases, consumption, features, risks/notes).
    Pass demo=true to include all historical meetings when the current week has none.
    """
    start, end = _week_bounds(week_start)
    records = _load_post_meetings(demo, start, end)

    if not records:
        return {
            "ok": True,
            "week_start": start.date().isoformat(),
            "week_end": (end - timedelta(days=1)).date().isoformat(),
            "slides": [],
            "companies": 0,
            "meetings": 0,
            "demo": demo,
        }

    grouped = _group_by_company(records)
    slides = []
    for company_name, company_records in grouped.items():
        sf = _extract_sf(company_records)
        if demo:
            slide = _mock_slide(company_name)
        else:
            try:
                slide = _build_slide(company_name, company_records, sf)
            except RuntimeError as exc:
                raise HTTPException(
                    status_code=503,
                    detail=f"Elastic inference connector required for customer data. {exc}",
                )
        slides.append(slide)

    return {
        "ok": True,
        "week_start": start.date().isoformat(),
        "week_end": (end - timedelta(days=1)).date().isoformat(),
        "slides": slides,
        "companies": len(slides),
        "meetings": len(records),
        "demo": demo,
    }
