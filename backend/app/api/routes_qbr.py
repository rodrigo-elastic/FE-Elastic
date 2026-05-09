"""
filename: routes_qbr.py
description: QBR (Quarterly Business Review) generator. AE-facing executive presentation
generator. Aggregates post-meeting records for a given quarter and company, uses Claude
to synthesize a 4-section QBR narrative (Look Back, Current State, Look Forward, Next Steps),
and exports a 4-slide PPTX deck.
date: 09-05-2026
"""
__author__ = "Rodrigo Careaga"
__copyright__ = "Copyright 2026, Rodrigo Careaga"
__version__ = "0.1.0"
__status__ = "Development"

import json
from datetime import date, datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel, Field

from app.config import settings
from app.integrations.claude_client import get_elastic_service
from app.utils.logging import get_logger

log = get_logger(__name__)

router = APIRouter(prefix="/qbr", tags=["qbr"])


# ============================================================ Pydantic models =


class QBRMetric(BaseModel):
    label: str
    value: str
    trend: str = ""  # "up", "down", "stable"


class QBRContent(BaseModel):
    company_name: str = ""
    quarter: str = ""
    use_case: str = ""
    arr: str = ""
    # Look Back
    kpis: List[QBRMetric] = Field(default_factory=list)
    technical_wins: List[str] = Field(default_factory=list)
    business_outcomes: List[str] = Field(default_factory=list)
    # Current State
    health_score: int = 85
    health_summary: str = ""
    feature_gaps: List[str] = Field(default_factory=list)
    optimization_recs: List[str] = Field(default_factory=list)
    # Look Forward
    expansion_opportunities: List[str] = Field(default_factory=list)
    new_use_cases: List[str] = Field(default_factory=list)
    roadmap_items: List[str] = Field(default_factory=list)
    # Next Steps
    next_steps: List[str] = Field(default_factory=list)


_QBR_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "company_name": {"type": "string"},
        "quarter": {"type": "string"},
        "use_case": {"type": "string"},
        "arr": {"type": "string"},
        "kpis": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "label": {"type": "string"},
                    "value": {"type": "string"},
                    "trend": {"type": "string"},
                },
                "required": ["label", "value"],
            },
        },
        "technical_wins": {"type": "array", "items": {"type": "string"}},
        "business_outcomes": {"type": "array", "items": {"type": "string"}},
        "health_score": {"type": "integer"},
        "health_summary": {"type": "string"},
        "feature_gaps": {"type": "array", "items": {"type": "string"}},
        "optimization_recs": {"type": "array", "items": {"type": "string"}},
        "expansion_opportunities": {"type": "array", "items": {"type": "string"}},
        "new_use_cases": {"type": "array", "items": {"type": "string"}},
        "roadmap_items": {"type": "array", "items": {"type": "string"}},
        "next_steps": {"type": "array", "items": {"type": "string"}},
    },
    "required": [
        "company_name", "quarter", "use_case", "kpis",
        "technical_wins", "business_outcomes", "health_score",
        "health_summary", "feature_gaps", "optimization_recs",
        "expansion_opportunities", "new_use_cases", "roadmap_items", "next_steps",
    ],
}


# ============================================================ Helpers ==============


def _current_quarter() -> str:
    """Return current quarter as 'Q2 2026' style string."""
    today = date.today()
    q = (today.month - 1) // 3 + 1
    return f"Q{q} {today.year}"


def _quarter_bounds(quarter_str: str) -> Tuple[datetime, datetime]:
    """Parse 'Q2 2026' into (start, end) datetime bounds (UTC)."""
    parts = quarter_str.strip().split()
    if len(parts) != 2 or not parts[0].upper().startswith("Q"):
        raise HTTPException(status_code=400, detail=f"Invalid quarter format: '{quarter_str}'. Use 'Q2 2026'.")
    try:
        q = int(parts[0][1:])
        year = int(parts[1])
    except ValueError:
        raise HTTPException(status_code=400, detail=f"Invalid quarter format: '{quarter_str}'. Use 'Q2 2026'.")
    if q < 1 or q > 4:
        raise HTTPException(status_code=400, detail=f"Quarter must be 1-4, got {q}.")
    month_start = (q - 1) * 3 + 1
    start = datetime(year, month_start, 1, tzinfo=timezone.utc)
    # End: first day of next quarter (or next year's Q1)
    if q == 4:
        end = datetime(year + 1, 1, 1, tzinfo=timezone.utc)
    else:
        end = datetime(year, month_start + 3, 1, tzinfo=timezone.utc)
    return start, end


def _load_meetings_for_account(
    company_id: str, start: datetime, end: datetime
) -> List[Dict[str, Any]]:
    """Scan runtime/post_meeting/*.json, filter by company and date range."""
    records: List[Dict[str, Any]] = []
    all_records: List[Dict[str, Any]] = []
    seen: set = set()

    post_dir = settings.runtime_dir / "post_meeting"
    if not post_dir.exists():
        return []

    for p in sorted(post_dir.glob("*.json")):
        try:
            rec = json.loads(p.read_text(encoding="utf-8"))
            mid = rec.get("meeting_id", p.stem)
            if mid in seen:
                continue

            # Match by company_id or company_name (case-insensitive contains)
            rec_company = (rec.get("company_name") or rec.get("company_id") or "").lower()
            if company_id.lower() not in rec_company and rec_company not in company_id.lower():
                continue

            seen.add(mid)
            all_records.append(rec)

            # Date filter within quarter bounds
            ga = rec.get("generated_at", "")
            if not ga:
                continue
            try:
                ts = datetime.fromisoformat(ga.replace("Z", "+00:00"))
                if start <= ts < end:
                    records.append(rec)
            except Exception:
                pass
        except Exception:
            pass

    # If no records matched the date filter, return all for the company
    return records if records else all_records


def _mock_qbr(company_name: str, quarter: str) -> Dict[str, Any]:
    """Searchlight Capital Splunk displacement demo narrative."""
    return {
        "company_name": company_name or "Searchlight Capital",
        "quarter": quarter,
        "use_case": "Security (SIEM) + Observability",
        "arr": "$4,600,000",
        "kpis": [
            {"label": "MTTR Reduction", "value": "40%", "trend": "up"},
            {"label": "Security Events/Day", "value": "4.2M", "trend": "up"},
            {"label": "Splunk Licenses Retired", "value": "3 of 4", "trend": "up"},
        ],
        "technical_wins": [
            "Successfully migrated 3 of 4 Splunk environments to Elastic SIEM - on schedule",
            "Deployed Elastic Agent across 6,200 endpoints with Fleet management centralization",
            "Built 47 custom detection rules using ES|QL replacing legacy SPL queries",
            "Integrated Elastic with Palo Alto Cortex XDR for unified threat correlation",
            "Achieved sub-5-second alert latency on critical security events vs 90s baseline",
        ],
        "business_outcomes": [
            "40% reduction in mean-time-to-respond (MTTR) for P1 security incidents",
            "Projected $1.2M annual savings vs Splunk renewal (avoided 62% cost increase)",
            "SOC analyst capacity freed up by 2 FTEs through automated triage workflows",
        ],
        "health_score": 91,
        "health_summary": (
            "Strong deployment health with active SOC engagement. Remaining Splunk migration "
            "milestone on track for Q3. Expansion pipeline open for SIEM + SOAR consolidation."
        ),
        "feature_gaps": [
            "Elastic AI Assistant - not yet enabled (SOC team evaluating use cases)",
            "Attack Discovery - licensed but not configured in production",
            "Entity Analytics - underutilized despite UEBA requirement in original POC scope",
        ],
        "optimization_recs": [
            "Enable Elastic AI Assistant for SOC L1 triage to accelerate alert investigation",
            "Configure Entity Analytics for insider threat detection - aligns with Q3 audit requirements",
            "Consolidate remaining Splunk forwarders using Elastic Agent to cut data duplication",
            "Schedule Attack Discovery workshop with SOC lead before Q3 renewal window opens",
        ],
        "expansion_opportunities": [
            "SOAR integration: replace legacy XSOAR runbooks with Elastic workflows ($800K NE opportunity)",
            "Cloud workload monitoring: AWS + Azure coverage currently dark - 3 business units",
            "GRC/compliance dashboards: PCI DSS and SOX reporting currently manual - FE POC ready",
        ],
        "new_use_cases": [
            "Insider threat program: Entity Analytics + behavioral baselines across privileged accounts",
            "Cloud-native SIEM for new AWS us-east-2 environment going live Q3",
            "Executive security reporting: custom dashboards for CISO board-level visibility",
        ],
        "roadmap_items": [
            "8.14 release: ES|QL enhancements directly address Searchlight's custom aggregation requests",
            "Elastic AI Assistant GA: production-ready for SOC deployments - aligns with Q3 rollout plan",
            "SOAR playbook builder (preview): targeted for Q4 - replaces XSOAR automation dependency",
        ],
        "next_steps": [
            "Rodrigo Careaga: Schedule Elastic AI Assistant enablement session by May 30",
            "CISO (Sarah Chen): Approve Entity Analytics production rollout by Jun 15",
            "Elastic PS (Marcus): Deliver final Splunk migration runbook for Environment 4 by Jun 1",
            "AE (Rodrigo): Send Q3 renewal proposal with SIEM expansion SKU by Jun 20",
        ],
    }


def _build_qbr_content(
    company_name: str,
    quarter: str,
    records: List[Dict[str, Any]],
    demo: bool = False,
) -> Dict[str, Any]:
    """Call Claude structured or return mock. Falls back to mock on RuntimeError."""
    if demo or not records:
        return _mock_qbr(company_name, quarter)

    summaries = [r.get("summary", "") for r in records if r.get("summary")]
    all_actions: List[Dict[str, Any]] = []
    for rec in records:
        all_actions.extend(rec.get("action_items") or [])
    meddpicc: List[Dict[str, Any]] = []
    for rec in records:
        meddpicc.extend(rec.get("meddpicc_signals") or [])

    ai_lines = "\n".join(
        f"- {ai.get('title', '')} (owner: {ai.get('owner_name', 'TBD')}, "
        f"due: {ai.get('due_date', 'TBD')}, impact: {ai.get('impact', 'med')})"
        for ai in all_actions
    ) or "No action items recorded."

    meddpicc_lines = "\n".join(
        f"- [{m.get('category', '')}] {m.get('note', '') or m.get('quote', '')[:120]}"
        for m in meddpicc[:10]
    ) or "No MEDDPICC signals."

    summary_text = "\n\n---\n\n".join(summaries) or "No meeting summaries available."

    system = (
        "You are a Field Engineering QBR specialist at Elastic. "
        "You produce structured, executive-ready Quarterly Business Review content for AE presentations. "
        "Be specific - use real names, amounts, dates from the input data. "
        "Never invent facts not in the input. If data is missing, provide reasonable estimates or say 'Data not available'."
    )

    user = f"""Generate a complete QBR (Quarterly Business Review) for the Elastic AE team.

Company: {company_name}
Quarter: {quarter}
Meeting count: {len(records)}

Meeting summaries:
{summary_text}

Action items:
{ai_lines}

MEDDPICC signals:
{meddpicc_lines}

Return a JSON object with ALL of these exact keys. Be thorough - executives need substance:
{{
  "company_name": "{company_name}",
  "quarter": "{quarter}",
  "use_case": "Primary Elastic use cases (e.g. 'Security + Observability')",
  "arr": "ARR if known from data, else empty string",
  "kpis": [
    {{"label": "Metric name", "value": "Measured value", "trend": "up|down|stable"}},
    {{"label": "Metric name", "value": "Measured value", "trend": "up|down|stable"}},
    {{"label": "Metric name", "value": "Measured value", "trend": "up|down|stable"}}
  ],
  "technical_wins": ["3-6 specific technical achievements this quarter"],
  "business_outcomes": ["2-4 measurable business outcomes delivered"],
  "health_score": 85,
  "health_summary": "2-3 sentence account health narrative",
  "feature_gaps": ["3-5 Elastic features that are underused or not yet enabled"],
  "optimization_recs": ["3-5 specific optimization recommendations with rationale"],
  "expansion_opportunities": ["2-4 expansion opportunities with estimated value if possible"],
  "new_use_cases": ["2-4 new Elastic use cases relevant to this account"],
  "roadmap_items": ["2-4 upcoming Elastic roadmap items relevant to this customer"],
  "next_steps": ["3-5 next steps in format: Owner: Action by Date"]
}}

Rules:
- health_score: 0-100 integer. 80+ = healthy, 60-79 = needs attention, <60 = at risk.
- kpis: exactly 3 metrics that matter most to this customer this quarter.
- next_steps: always include owner name and target date.
- Respond with ONLY the JSON object. No markdown. No explanation."""

    mock_payload = _mock_qbr(company_name, quarter)
    model_name = settings.model_for("post_meeting")

    try:
        svc = get_elastic_service()
        result: QBRContent = svc.call_structured(
            system=system,
            user=user,
            schema=_QBR_SCHEMA,
            output_model=QBRContent,
            model=model_name,
            max_tokens=2000,
            effort="high",
            thinking_adaptive=True,
            cache_system=True,
            mock_payload=mock_payload,
            audit_meta={"agent": "qbr", "company": company_name, "quarter": quarter},
            strict=True,
        )
        return result.model_dump()
    except RuntimeError as exc:
        log.warning("qbr.elastic_required_falling_back_to_mock",
                    company=company_name, quarter=quarter, error=str(exc)[:200])
        return mock_payload
    except Exception as exc:
        log.warning("qbr.claude_failed_falling_back_to_mock",
                    company=company_name, quarter=quarter, error=str(exc)[:200])
        return mock_payload


# ============================================================ PPTX builder =========


def _build_qbr_pptx(content: Dict[str, Any]) -> Path:
    """Generate 4-slide QBR PPTX deck."""
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.dml.color import RGBColor  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]
        from pptx.util import Inches, Pt  # type: ignore[import]
    except ImportError as exc:
        raise RuntimeError("python-pptx required: pip install python-pptx") from exc

    # ── Palette ───────────────────────────────────────────────────
    RED   = RGBColor(0xE8, 0x4B, 0x37)
    NAVY  = RGBColor(0x0F, 0x2D, 0x5C)
    TEAL  = RGBColor(0x00, 0xB4, 0xA2)
    GOLD  = RGBColor(0xF1, 0xA7, 0x30)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
    DGRAY = RGBColor(0xCC, 0xCC, 0xCC)
    DTEXT = RGBColor(0x1A, 0x1A, 0x1A)
    MUTED = RGBColor(0x88, 0x88, 0x88)
    EBLUE = RGBColor(0x0B, 0x64, 0xDD)
    GREEN = RGBColor(0x3C, 0xB4, 0x4B)

    I = Inches

    def _make_slide(prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _rect(slide, l, t, w, h, bg, border=None):
        s = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
        s.fill.solid()
        s.fill.fore_color.rgb = bg
        if border:
            s.line.color.rgb = border
            s.line.width = Pt(0.75)
        else:
            s.line.fill.background()
        return s

    def _new_tf(slide, l, t, w, h, wrap=True):
        box = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        return tf

    def _para(tf, text, size=8.0, bold=False, italic=False,
              color=None, align=PP_ALIGN.LEFT, space_before=0, underline=False):
        if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(space_before)
        r = p.add_run()
        r.text = str(text)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.underline = underline
        r.font.color.rgb = color if color else DTEXT
        return p

    def _hdr_bar(slide, l, t, w, h, title, bg, size=11, underline=False):
        s = _rect(slide, l, t, w, h, bg)
        p = s.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.underline = underline
        r.font.color.rgb = WHITE

    def _mini_hdr(slide, l, t, w, h, title, bg):
        s = _rect(slide, l, t, w, h, bg)
        p = s.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = WHITE

    def _bullet_list(slide, items, l, t, w, h, size=8.0, max_items=None, numbered=False):
        if max_items:
            items = items[:max_items]
        tf = _new_tf(slide, l, t, w, h)
        if not items:
            _para(tf, "None recorded.", size=size, color=MUTED, italic=True)
            return
        for i, item in enumerate(items):
            prefix = f"{i+1}. " if numbered else "- "
            _para(tf, f"{prefix}{item}", size=size, space_before=(3 if i > 0 else 0))

    prs = Presentation()
    prs.slide_width  = I(13.33)
    prs.slide_height = I(7.5)

    co     = content.get("company_name") or "Account"
    qtr    = content.get("quarter") or _current_quarter()
    today  = date.today().strftime("%B %d, %Y")

    # ==================================================================
    # SLIDE 1 - Cover
    # ==================================================================
    sl1 = _make_slide(prs)
    sl1.background.fill.solid()
    sl1.background.fill.fore_color.rgb = NAVY

    # Left accent bar
    _rect(sl1, 0, 0, 0.18, 7.5, RED)

    # elastic wordmark
    tf = _new_tf(sl1, 0.3, 0.2, 2.5, 0.36, wrap=False)
    _para(tf, "elastic", size=14, italic=True, color=WHITE)

    # Company name
    tf = _new_tf(sl1, 0.3, 2.8, 12.73, 0.9)
    _para(tf, co, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # QBR subtitle
    tf = _new_tf(sl1, 0.3, 3.75, 12.73, 0.45)
    _para(tf, "Quarterly Business Review", size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # Quarter
    tf = _new_tf(sl1, 0.3, 4.28, 12.73, 0.36)
    _para(tf, qtr, size=14, bold=True, color=EBLUE, align=PP_ALIGN.CENTER)

    # Updated date
    tf = _new_tf(sl1, 0.3, 4.72, 12.73, 0.30)
    _para(tf, f"Updated: {today}", size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # FE Copilot watermark bottom-right
    tf = _new_tf(sl1, 11.5, 7.1, 1.7, 0.3, wrap=False)
    _para(tf, "FE Copilot", size=9, color=WHITE, align=PP_ALIGN.RIGHT)

    # ==================================================================
    # SLIDE 2 - Look Back: Value Delivered
    # ==================================================================
    sl2 = _make_slide(prs)
    sl2.background.fill.solid()
    sl2.background.fill.fore_color.rgb = WHITE

    _hdr_bar(sl2, 0, 0, 13.33, 0.38, "Look Back: Value Delivered", RED, size=11)

    # KPI boxes - 3 side by side
    kpis = content.get("kpis") or []
    kpi_w = 4.2
    kpi_gap = (13.33 - 0.2 - kpi_w * 3) / 2
    kpi_y = 0.50
    kpi_h = 1.5
    for i, kpi in enumerate(kpis[:3]):
        kx = 0.1 + i * (kpi_w + kpi_gap)
        _rect(sl2, kx, kpi_y, kpi_w, kpi_h, LGRAY, border=DGRAY)
        # Label
        tf = _new_tf(sl2, kx + 0.1, kpi_y + 0.1, kpi_w - 0.2, 0.28)
        _para(tf, kpi.get("label", ""), size=8, color=MUTED, align=PP_ALIGN.CENTER)
        # Value
        tf = _new_tf(sl2, kx + 0.1, kpi_y + 0.4, kpi_w - 0.2, 0.7)
        _para(tf, kpi.get("value", ""), size=28, bold=True, align=PP_ALIGN.CENTER)
        # Trend arrow
        trend = kpi.get("trend", "")
        if trend == "up":
            arrow, arrow_color = "▲", GREEN
        elif trend == "down":
            arrow, arrow_color = "▼", RED
        else:
            arrow, arrow_color = "→", GOLD
        tf = _new_tf(sl2, kx + 0.1, kpi_y + 1.15, kpi_w - 0.2, 0.28)
        _para(tf, arrow, size=11, color=arrow_color, align=PP_ALIGN.CENTER)

    # Technical Wins
    tf = _new_tf(sl2, 0.1, 2.15, 8.4, 0.28)
    _para(tf, "Technical Wins", size=10, bold=True, color=RED)

    wins = content.get("technical_wins") or []
    _bullet_list(sl2, wins, 0.1, 2.50, 8.4, 2.2, size=8.5, max_items=5)

    # Business Outcomes
    tf = _new_tf(sl2, 0.1, 4.75, 8.4, 0.28)
    _para(tf, "Business Outcomes", size=10, bold=True, color=NAVY)

    outcomes = content.get("business_outcomes") or []
    _bullet_list(sl2, outcomes, 0.1, 5.08, 8.4, 1.5, size=8.5, max_items=3)

    # Right side panel - Quarter at a Glance
    PANEL_X = 8.8
    PANEL_W = 4.43
    _rect(sl2, PANEL_X, 0.5, PANEL_W, 6.8, LGRAY)
    tf = _new_tf(sl2, PANEL_X + 0.12, 0.62, PANEL_W - 0.24, 0.32)
    _para(tf, "Quarter at a Glance", size=9, bold=True, color=NAVY)

    # ARR pill
    arr_val = content.get("arr") or "-"
    _rect(sl2, PANEL_X + 0.12, 1.05, PANEL_W - 0.24, 0.36, DGRAY)
    tf = _new_tf(sl2, PANEL_X + 0.2, 1.08, PANEL_W - 0.4, 0.30)
    _para(tf, f"ARR: {arr_val}", size=9, bold=True)

    # Use case pill
    use_case_val = content.get("use_case") or "-"
    _rect(sl2, PANEL_X + 0.12, 1.50, PANEL_W - 0.24, 0.36, DGRAY)
    tf = _new_tf(sl2, PANEL_X + 0.2, 1.53, PANEL_W - 0.4, 0.30)
    _para(tf, f"Use Case: {use_case_val}", size=8.5)

    # Health score pill
    hs = content.get("health_score", 85)
    hs_color = GREEN if hs >= 80 else (GOLD if hs >= 60 else RED)
    _rect(sl2, PANEL_X + 0.12, 1.95, PANEL_W - 0.24, 0.36, DGRAY)
    tf = _new_tf(sl2, PANEL_X + 0.2, 1.98, PANEL_W - 0.4, 0.30)
    _para(tf, f"Health Score: {hs}/100", size=9, bold=True, color=hs_color)

    # Quarter label
    tf = _new_tf(sl2, PANEL_X + 0.12, 2.45, PANEL_W - 0.24, 0.28)
    _para(tf, f"Quarter: {qtr}", size=8.5, color=MUTED)

    # ==================================================================
    # SLIDE 3 - Current State: Deployment Health
    # ==================================================================
    sl3 = _make_slide(prs)
    sl3.background.fill.solid()
    sl3.background.fill.fore_color.rgb = WHITE

    _hdr_bar(sl3, 0, 0, 13.33, 0.38, "Current State: Deployment Health", NAVY, size=11)

    # Health score circle (rounded rect)
    hs_val = content.get("health_score", 85)
    hs_color = GREEN if hs_val >= 80 else (GOLD if hs_val >= 60 else RED)
    _rect(sl3, 0.3, 0.6, 2.5, 2.5, LGRAY)

    # Big number
    tf = _new_tf(sl3, 0.3, 0.75, 2.5, 1.4)
    _para(tf, str(hs_val), size=48, bold=True, color=hs_color, align=PP_ALIGN.CENTER)

    # /100
    tf = _new_tf(sl3, 0.3, 2.22, 2.5, 0.35)
    _para(tf, "/100", size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # Health Score label
    tf = _new_tf(sl3, 0.3, 2.60, 2.5, 0.32)
    _para(tf, "Health Score", size=9, color=DTEXT, align=PP_ALIGN.CENTER)

    # Health summary
    hs_summary = content.get("health_summary", "")
    tf = _new_tf(sl3, 0.3, 3.2, 2.8, 1.0)
    _para(tf, hs_summary, size=8, italic=True, color=MUTED)

    # Feature Gaps section
    GAPS_X = 3.2
    GAPS_W = 4.5
    _mini_hdr(sl3, GAPS_X, 0.6, GAPS_W, 0.28, "Feature Gaps", GOLD)
    _rect(sl3, GAPS_X, 0.88, GAPS_W, 3.6, LGRAY)
    gaps = content.get("feature_gaps") or []
    _bullet_list(sl3, gaps, GAPS_X + 0.1, 0.98, GAPS_W - 0.2, 3.4, size=8.0, max_items=6)

    # Optimization Recommendations section
    OPT_X = 8.0
    OPT_W = 5.2
    _mini_hdr(sl3, OPT_X, 0.6, OPT_W, 0.28, "Optimization Recommendations", RED)
    _rect(sl3, OPT_X, 0.88, OPT_W, 3.6, LGRAY)
    opt_recs = content.get("optimization_recs") or []
    _bullet_list(sl3, opt_recs, OPT_X + 0.1, 0.98, OPT_W - 0.2, 3.4, size=8.0, max_items=5, numbered=True)

    # ==================================================================
    # SLIDE 4 - Look Forward: Strategic Roadmap
    # ==================================================================
    sl4 = _make_slide(prs)
    sl4.background.fill.solid()
    sl4.background.fill.fore_color.rgb = WHITE

    _hdr_bar(sl4, 0, 0, 13.33, 0.38, "Look Forward: Strategic Roadmap", TEAL, size=11, underline=True)

    # Three columns
    COL_W = 4.3
    COL_GAP = 0.115
    COLS = [
        (0.0 + COL_GAP, "Expansion Opportunities", GOLD, "expansion_opportunities"),
        (COL_W + COL_GAP * 2, "New Use Cases", NAVY, "new_use_cases"),
        (COL_W * 2 + COL_GAP * 3, "Roadmap Alignment", TEAL, "roadmap_items"),
    ]
    COL_TOP = 0.50
    COL_BODY_H = 5.5

    for col_x, col_title, col_color, col_key in COLS:
        _mini_hdr(sl4, col_x, COL_TOP, COL_W, 0.30, col_title, col_color)
        _rect(sl4, col_x, COL_TOP + 0.30, COL_W, COL_BODY_H, LGRAY)
        items = content.get(col_key) or []
        _bullet_list(sl4, items, col_x + 0.1, COL_TOP + 0.40, COL_W - 0.2,
                     COL_BODY_H - 0.2, size=8.0, max_items=4)

    # Bottom Next Steps strip
    NS_Y = 6.2
    NS_H = 1.2
    _rect(sl4, 0, NS_Y, 13.33, NS_H, LGRAY)
    tf = _new_tf(sl4, 0.15, NS_Y + 0.08, 1.6, 0.32)
    _para(tf, "Next Steps", size=9, bold=True, color=RED)

    next_steps = content.get("next_steps") or []
    if next_steps:
        steps_text = "  |  ".join(next_steps[:4])
    else:
        steps_text = "No next steps defined."
    tf = _new_tf(sl4, 0.15, NS_Y + 0.42, 13.0, 0.7)
    _para(tf, steps_text, size=8, color=DTEXT)

    # ── Save ─────────────────────────────────────────────────────
    slides_dir = settings.runtime_dir / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)
    slug = "".join(c if c.isalnum() or c in "-_" else "_"
                   for c in (content.get("company_name") or "Account"))
    qtr_slug = (content.get("quarter") or _current_quarter()).replace(" ", "_")
    out = slides_dir / f"{slug}_QBR_{qtr_slug}.pptx"
    prs.save(str(out))
    log.info("qbr.pptx_saved", path=str(out))
    return out


# ============================================================ Endpoints =============


class _GenerateRequest(BaseModel):
    company_id: str
    quarter: Optional[str] = None
    demo: bool = False


@router.get("/accounts")
def list_accounts() -> Dict[str, Any]:
    """List unique company IDs derived from post-meeting filenames."""
    post_dir = settings.runtime_dir / "post_meeting"
    companies: List[str] = []
    seen: set = set()

    if post_dir.exists():
        for p in sorted(post_dir.glob("*.json")):
            try:
                rec = json.loads(p.read_text(encoding="utf-8"))
                name = (rec.get("company_name") or rec.get("company_id") or "").strip()
                if name and name not in seen:
                    seen.add(name)
                    companies.append(name)
            except Exception:
                pass

    return {"ok": True, "accounts": companies, "count": len(companies)}


@router.post("/generate")
def generate_qbr(body: _GenerateRequest) -> Dict[str, Any]:
    """Generate a QBR for the given company and quarter.

    Uses the Elastic inference connector (strict mode - customer data never
    leaves Elastic infrastructure). Pass demo=true to use mock/Searchlight content.
    Falls back to mock on RuntimeError.
    """
    company_id = (body.company_id or "").strip()
    if not company_id:
        raise HTTPException(status_code=400, detail="company_id is required.")

    quarter = (body.quarter or _current_quarter()).strip()

    # Validate quarter format early
    try:
        start, end = _quarter_bounds(quarter)
    except HTTPException:
        raise

    # Load meeting records for this account + quarter
    records = _load_meetings_for_account(company_id, start, end)

    # Build QBR content
    content = _build_qbr_content(
        company_name=company_id,
        quarter=quarter,
        records=records,
        demo=body.demo,
    )

    # Build PPTX
    try:
        pptx_path = _build_qbr_pptx(content)
    except Exception as exc:
        log.warning("qbr.pptx_failed", company=company_id, error=str(exc)[:200])
        raise HTTPException(status_code=500, detail=f"PPTX generation failed: {exc}")

    pptx_rel = f"/api/v1/qbr/download/{pptx_path.name}"
    pptx_url = f"{settings.public_base_url}{pptx_rel}"

    log.info("qbr.generate_done",
             company=company_id, quarter=quarter,
             meetings=len(records), demo=body.demo,
             pptx=pptx_path.name)

    return {
        "ok": True,
        "company_name": content.get("company_name", company_id),
        "quarter": quarter,
        "slide_name": pptx_path.name,
        "pptx_url": pptx_url,
        "pptx_rel": pptx_rel,
        "content": content,
    }


@router.get("/download/{filename}")
def download_qbr(filename: str):
    """Serve a generated QBR PPTX file."""
    from fastapi.responses import FileResponse

    safe = "".join(c for c in filename if c.isalnum() or c in "-_.")
    path = settings.runtime_dir / "slides" / safe
    if not path.exists() or path.suffix != ".pptx":
        raise HTTPException(status_code=404, detail="QBR slide not found.")
    return FileResponse(
        str(path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=safe,
    )
